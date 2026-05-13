"""
AlphaPULSE Slide Generator
Reads the research brief (.md) and populates the AlphaPULSE PPT template.

Run: python scripts/generate_alphapulse.py
Output: output/AlphaPULSE_YYMMDD_draft.pptx

The script:
  1. Copies the latest AlphaPULSE template
  2. Updates cover slide: date + weekly theme title (you enter these)
  3. Populates a new "Key Factors" text slide with the brief content
  4. Leaves Bloomberg chart placeholder slots labeled for manual insertion
  5. Preserves disclaimer slides verbatim
"""

import sys
import os
if sys.stdout.encoding != "utf-8":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

import re
import glob
import shutil
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# ── Paths ───────────────────────────────────────────────────────────────────────
ROOT         = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE_DIR = os.path.join(ROOT, "template")
OUTPUT_DIR   = os.path.join(ROOT, "output")
DATA_DIR     = os.path.join(ROOT, "data")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# AlphaPULSE brand colours  (#1F3864 = dark navy, #C9A02B = gold)
DARK_BLUE  = RGBColor(0x1F, 0x38, 0x64)
GOLD       = RGBColor(0xC9, 0xA0, 0x2B)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY  = RGBColor(0x40, 0x40, 0x40)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)

# ── Helpers ─────────────────────────────────────────────────────────────────────

def find_latest_template() -> str:
    """Return the most recently modified AlphaPULSE template."""
    files = sorted(
        glob.glob(os.path.join(TEMPLATE_DIR, "WealthStrategy_AlphaPULSE_*.pptx")),
        key=os.path.getmtime,
        reverse=True,
    )
    if not files:
        raise FileNotFoundError(f"No AlphaPULSE template found in {TEMPLATE_DIR}")
    return files[0]


def find_latest_brief() -> str:
    """Return the most recently generated research brief."""
    files = sorted(
        glob.glob(os.path.join(OUTPUT_DIR, "research_brief_*.md")),
        key=os.path.getmtime,
        reverse=True,
    )
    if not files:
        raise FileNotFoundError(
            f"No research brief found in {OUTPUT_DIR}. Run research_brief.py first."
        )
    return files[0]


def parse_brief(path: str) -> dict:
    """Parse sections from the markdown research brief."""
    with open(path, encoding="utf-8") as f:
        text = f.read()

    # Strip YAML front matter
    text = re.sub(r"^---.*?---\n", "", text, flags=re.DOTALL)

    sections = {}
    current_section = None
    buffer = []

    for line in text.splitlines():
        if line.startswith("## "):
            if current_section:
                sections[current_section] = "\n".join(buffer).strip()
            current_section = line[3:].strip()
            buffer = []
        elif line.startswith("# "):
            sections["title"] = line[2:].strip()
        else:
            buffer.append(line)

    if current_section:
        sections[current_section] = "\n".join(buffer).strip()

    return sections


def set_text_run(run, text, bold=False, size_pt=None, color=None):
    run.text = text
    run.font.bold = bold
    if size_pt:
        run.font.size = Pt(size_pt)
    if color:
        run.font.color.rgb = color


def add_text_box(slide, left, top, width, height, text,
                 font_size=11, bold=False, color=DARK_GREY,
                 bg_color=None, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if bg_color:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return txBox


def add_placeholder_image_box(slide, left, top, width, height, label="INSERT CHART FROM BLOOMBERG"):
    """Add a visible placeholder box where Bloomberg charts should be inserted."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_GREY
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[CHART] {label}"
    run.font.size = Pt(10)
    run.font.color.rgb = DARK_BLUE
    run.font.bold = True
    return shape

# ── Slide builders ──────────────────────────────────────────────────────────────

def build_factors_slide(prs: Presentation, sections: dict, theme_title: str) -> None:
    """Add a new Key Factors content slide."""
    # Use blank layout (last layout usually blank)
    blank_layout = prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    # Header bar
    header = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.33), Inches(0.6)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()

    # Header text
    add_text_box(slide, 0.15, 0.1, 10, 0.45,
                 f"AlphaPULSE | ปัจจัยขับเคลื่อนตลาด — {theme_title}",
                 font_size=12, bold=True, color=WHITE)

    # Macro regime summary
    macro_text = sections.get("ภาพรวม Macro Regime", "")
    if macro_text:
        add_text_box(slide, 0.3, 0.7, 12.7, 0.8, macro_text,
                     font_size=9.5, color=DARK_GREY)

    # Key factors text
    factors_text = sections.get("ปัจจัยขับเคลื่อนตลาด (Key Factors)", "")
    if factors_text:
        add_text_box(slide, 0.3, 1.55, 12.7, 4.8, factors_text,
                     font_size=9.5, color=DARK_GREY)

    # Strategy box
    strategy_text = sections.get("กลยุทธ์การลงทุน", "")
    if strategy_text:
        add_text_box(slide, 0.3, 6.4, 6.2, 1.0,
                     ">> กลยุทธ์การลงทุน",
                     font_size=10, bold=True, color=WHITE, bg_color=DARK_BLUE)
        add_text_box(slide, 0.3, 7.0, 6.2, 0.9, strategy_text,
                     font_size=9, color=DARK_GREY, bg_color=LIGHT_GREY)

    # Risk box
    risk_text = sections.get("ความเสี่ยงที่ต้องติดตาม", "")
    if risk_text:
        add_text_box(slide, 6.8, 6.4, 6.2, 1.0,
                     "!! ความเสี่ยงที่ต้องติดตาม",
                     font_size=10, bold=True, color=WHITE, bg_color=RGBColor(0xC0, 0x39, 0x2B))
        add_text_box(slide, 6.8, 7.0, 6.2, 0.9, risk_text,
                     font_size=9, color=DARK_GREY, bg_color=LIGHT_GREY)


def build_chart_slide(prs: Presentation, theme_title: str) -> None:
    """Add a chart placeholder slide for Bloomberg charts."""
    blank_layout = prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    # Header bar
    header = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.33), Inches(0.6)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.fill.background()

    add_text_box(slide, 0.15, 0.1, 10, 0.45,
                 f"AlphaPULSE | สนับสนุนการวิเคราะห์ — {theme_title}",
                 font_size=12, bold=True, color=WHITE)

    # Two chart placeholders side by side
    add_placeholder_image_box(slide, 0.3, 0.8, 6.1, 3.5,
                               "INSERT CHART 1 FROM BLOOMBERG\n(e.g. SET Index + Foreign Flow)")
    add_placeholder_image_box(slide, 6.7, 0.8, 6.1, 3.5,
                               "INSERT CHART 2 FROM BLOOMBERG\n(e.g. Earnings Revision Breadth)")

    # Bottom notes area
    add_text_box(slide, 0.3, 4.5, 12.7, 0.4,
                 "วิธีใช้: เปิด Bloomberg/Excel → Copy chart → Paste ใน PowerPoint แทนที่กล่องสีเทาข้างต้น",
                 font_size=8.5, color=DARK_GREY)

    # Sector table placeholder
    add_placeholder_image_box(slide, 0.3, 5.0, 12.7, 2.3,
                               "INSERT SECTOR PERFORMANCE TABLE FROM BLOOMBERG / BLS SYSTEM")

# ── Cover slide update ──────────────────────────────────────────────────────────

def update_cover_slide(prs: Presentation, date_str: str, theme_title: str,
                       top_bullet: str) -> None:
    """Try to update date and theme on slide 0 (cover). Non-destructive."""
    slide = prs.slides[0]
    updated = {"date": False, "theme": False}

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full_text = "".join(r.text for r in para.runs).strip()

            # Look for existing date pattern (e.g. "May 11, 2026" or "May 2026")
            if re.search(r"\b(January|February|March|April|May|June|July|August|"
                         r"September|October|November|December)\b.*\d{4}", full_text):
                if not updated["date"]:
                    for run in para.runs:
                        if re.search(r"\d{4}", run.text):
                            run.text = date_str
                            updated["date"] = True
                            break

    # Add date + theme text box on cover if we couldn't find existing fields
    if not updated["date"]:
        add_text_box(slide, 0.5, 6.5, 9, 0.5, date_str,
                     font_size=13, bold=True, color=DARK_BLUE)

    # Always add theme as overlay (non-destructive — doesn't remove existing shapes)
    add_text_box(slide, 0.5, 7.1, 12, 0.6, f"Theme: {theme_title}",
                 font_size=12, color=DARK_GREY)

    if top_bullet:
        add_text_box(slide, 0.5, 7.75, 12, 0.8, f"• {top_bullet}",
                     font_size=10, color=DARK_GREY)

# ── Main ─────────────────────────────────────────────────────────────────────────

def main():
    today    = datetime.today()
    date_str = today.strftime("%d %B %Y")
    date_file = today.strftime("%y%m%d")

    print("=" * 60)
    print("AlphaPULSE Slide Generator")
    print(f"Date: {date_str}")
    print("=" * 60)

    # User inputs
    print("\nEnter the weekly theme title (English or Thai):")
    print("  Example: '1Q26 Earnings Checkpoint' or 'Reality Check Phase'")
    theme_title = input("  Theme: ").strip()
    if not theme_title:
        theme_title = f"Weekly Strategy — {date_str}"

    # Find files
    template_path = find_latest_template()
    brief_path    = find_latest_brief()
    print(f"\nTemplate:      {os.path.basename(template_path)}")
    print(f"Research brief: {os.path.basename(brief_path)}")

    # Parse brief
    print("\nParsing research brief...")
    sections = parse_brief(brief_path)
    top_bullet = ""
    factors_text = sections.get("ปัจจัยขับเคลื่อนตลาด (Key Factors)", "")
    if factors_text:
        # Extract first factor as the cover bullet
        first_line = [l.strip() for l in factors_text.splitlines() if l.strip()]
        top_bullet = first_line[0][:120] if first_line else ""

    # Copy template → output
    output_path = os.path.join(OUTPUT_DIR, f"AlphaPULSE_{date_file}_draft.pptx")
    shutil.copy2(template_path, output_path)
    print(f"Copied template to: {os.path.basename(output_path)}")

    # Load and modify
    prs = Presentation(output_path)

    # Update cover slide
    print("Updating cover slide...")
    update_cover_slide(prs, date_str, theme_title, top_bullet)

    # Insert new slides BEFORE the disclaimer slides (last 2 slides)
    # Strategy: add new slides to end, then we rely on user to reorder if needed
    print("Adding Key Factors slide...")
    build_factors_slide(prs, sections, theme_title)

    print("Adding Bloomberg chart placeholder slide...")
    build_chart_slide(prs, theme_title)

    # Save
    prs.save(output_path)
    print(f"\nDraft saved: {output_path}")

    print("\n── NEXT STEPS ───────────────────────────────────────────")
    print("1. Open the PPTX file in PowerPoint")
    print("2. Move the new slides to the correct position (after cover)")
    print("3. Replace grey placeholder boxes with Bloomberg chart images")
    print("4. Review and adjust Thai text as needed")
    print("5. The last 2 slides (Disclaimers) are unchanged from template")
    print("─" * 60)


if __name__ == "__main__":
    main()
